"""LangGraph API 구조 계획 PPT 생성 스크립트."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 색상 팔레트 ──────────────────────────────────
BG_DARK = RGBColor(0x1E, 0x1E, 0x2E)       # 진한 남색 배경
BG_SECTION = RGBColor(0x28, 0x28, 0x3C)     # 섹션 배경
ACCENT = RGBColor(0x89, 0xB4, 0xFA)         # 밝은 파랑 강조
ACCENT2 = RGBColor(0xA6, 0xE3, 0xA1)        # 민트 그린
ACCENT3 = RGBColor(0xF9, 0xE2, 0xAF)        # 따뜻한 노랑
WHITE = RGBColor(0xF5, 0xF5, 0xF5)
GRAY = RGBColor(0xA0, 0xA0, 0xB0)
ORANGE = RGBColor(0xFA, 0xB3, 0x87)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, color=WHITE, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "맑은 고딕"
        p.space_after = Pt(6)
        p.level = 0
    return tf


def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.6), color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rounded_box(slide, left, top, width, height, text,
                    fill_color=BG_SECTION, text_color=WHITE, font_size=14, bold=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = "맑은 고딕"
    p.font.bold = bold
    return shape


def section_header(slide, number, title):
    add_accent_bar(slide, Inches(0.6), Inches(0.5), Inches(0.08), Inches(0.55), ACCENT)
    add_textbox(slide, Inches(0.85), Inches(0.45), Inches(10), Inches(0.7),
                f"{number}. {title}", font_size=32, color=ACCENT, bold=True)
    # subtle divider line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.15), Inches(12), Pt(1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x40, 0x40, 0x55)
    shape.line.fill.background()


# ══════════════════════════════════════════════════
# Slide 1: Title
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
            "LangGraph 기반 API 구조 계획",
            font_size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
            "LLM Chatbot 프로젝트 아키텍처 설계 문서",
            font_size=22, color=GRAY, align=PP_ALIGN.CENTER)

# decorative line
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.2), Inches(4.3), Pt(3)
)
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_textbox(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.6),
            "Workflow Orchestration  ·  MCP Tool Calling  ·  RAG  ·  Agent",
            font_size=18, color=ACCENT2, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════
# Slide 2: 목적
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "01", "목적")

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.7),
            "현재: Cube → queue → worker → LLM → Cube 응답 (단순 흐름)",
            font_size=18, color=GRAY)

add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(0.6),
            "앞으로 수용해야 할 방향:",
            font_size=20, color=WHITE, bold=True)

items = [
    "LangGraph 기반 workflow orchestration",
    "Conditional routing",
    "MCP tool calling",
    "RAG workflow 추가",
    "Agent workflow 추가",
]
add_bullet_list(slide, Inches(1.0), Inches(2.9), Inches(10), Inches(3.0),
                items, font_size=20, color=WHITE)

add_rounded_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.0),
                "이 문서는 위 요구를 반영해 api/ 폴더를 어떻게 나누는 것이 적절한지 정리한다.",
                fill_color=RGBColor(0x2A, 0x2A, 0x45), text_color=ACCENT3, font_size=18)

# ══════════════════════════════════════════════════
# Slide 3: 핵심 설계 원칙 (1/2)
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "02", "핵심 설계 원칙 (1/3)")

# Principle 1
add_accent_bar(slide, Inches(0.8), Inches(1.5), Inches(0.06), Inches(1.6), ACCENT)
add_textbox(slide, Inches(1.1), Inches(1.45), Inches(5.3), Inches(0.5),
            "유스케이스와 workflow 구현을 분리", font_size=20, color=ACCENT2, bold=True)
add_textbox(slide, Inches(1.1), Inches(1.95), Inches(5.3), Inches(1.2),
            "api/chat/ → 애플리케이션 계층 (유스케이스)\n"
            "api/workflows/chat/ → LangGraph 구현 계층\n"
            "→ workflow 엔진을 바꿔도 진입점 유지 가능",
            font_size=15, color=WHITE)

# Principle 2
add_accent_bar(slide, Inches(6.8), Inches(1.5), Inches(0.06), Inches(1.6), ACCENT)
add_textbox(slide, Inches(7.1), Inches(1.45), Inches(5.5), Inches(0.5),
            "MCP는 처음부터 확장형 구조", font_size=20, color=ACCENT2, bold=True)
add_textbox(slide, Inches(7.1), Inches(1.95), Inches(5.5), Inches(1.2),
            "api/mcp/ → 별도 인프라 계층\n"
            "서버 레지스트리, tool 캐시, schema 변환, 실행 라우팅\n"
            "→ MCP 서버 수 증가에 대비",
            font_size=15, color=WHITE)

# Principle 3
add_accent_bar(slide, Inches(0.8), Inches(3.6), Inches(0.06), Inches(2.0), ACCENT)
add_textbox(slide, Inches(1.1), Inches(3.55), Inches(5.3), Inches(0.5),
            "RAG와 Agent는 점진적 승격", font_size=20, color=ACCENT2, bold=True)
add_textbox(slide, Inches(1.1), Inches(4.1), Inches(5.3), Inches(1.5),
            "초기: chat workflow 안에서 node 수준 호출\n"
            "복잡해지면 workflows/rag, workflows/agent로 분리\n\n"
            "• 단순 retrieval → api/rag/\n"
            "• 독립 RAG subgraph → api/workflows/rag/\n"
            "• 독립 agent loop → api/workflows/agent/",
            font_size=15, color=WHITE)

# Principle 4
add_accent_bar(slide, Inches(6.8), Inches(3.6), Inches(0.06), Inches(2.0), ACCENT)
add_textbox(slide, Inches(7.1), Inches(3.55), Inches(5.5), Inches(0.5),
            "Cube 계층은 얇게 유지", font_size=20, color=ACCENT2, bold=True)
add_textbox(slide, Inches(7.1), Inches(4.1), Inches(5.5), Inches(1.5),
            "cube/service.py → 응답 포맷 라우팅만 담당\n"
            "chat/service.py → ChatResult 구조체 반환\n"
            "chat/ 계층은 Cube를 전혀 모름\n\n"
            "→ 테스트·재사용이 쉬운 구조",
            font_size=15, color=WHITE)

# ══════════════════════════════════════════════════
# Slide 4: 핵심 설계 원칙 (2/3) - ChatResult
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "02", "핵심 설계 원칙 (2/3)")

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
            "ChatResult 구조체", font_size=22, color=ACCENT2, bold=True)

code_text = (
    "@dataclass(frozen=True, slots=True)\n"
    "class ChatResult:\n"
    "    reply: str            # LLM 텍스트 답변\n"
    "    model_used: str       # 실제 사용 모델\n"
    "    tool_calls: list[str] # 호출된 tool 목록\n"
    "    thread_id: str        # checkpointer thread_id"
)
add_rounded_box(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.5),
                code_text, fill_color=RGBColor(0x20, 0x20, 0x35),
                text_color=ACCENT3, font_size=15)

reasons = [
    "cube/service.py가 응답 포맷 선택에 메타데이터 필요",
    "archive/service.py가 model_used, tool_calls 인덱싱",
    "reply만 반환하면 계층 분리가 무너짐",
]
add_textbox(slide, Inches(0.8), Inches(4.8), Inches(5.5), Inches(0.4),
            "반환 이유:", font_size=16, color=ACCENT, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(5.2), Inches(5.5), Inches(1.5),
                reasons, font_size=14, color=WHITE)

# Right side - Cube 응답 포맷
add_textbox(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.5),
            "Cube 응답 포맷 전략", font_size=22, color=ACCENT2, bold=True)

# multimessage box
add_rounded_box(slide, Inches(6.8), Inches(2.2), Inches(2.7), Inches(2.0),
                "multimessage\n\n순수 텍스트\n복사 가능\n현재 사용 중",
                fill_color=RGBColor(0x25, 0x35, 0x25), text_color=ACCENT2, font_size=14)

# richnotification box
add_rounded_box(slide, Inches(9.8), Inches(2.2), Inches(2.7), Inches(2.0),
                "richnotification\n\n이미지 기반 렌더링\n테이블·버튼 지원\n이후 추가 예정",
                fill_color=RGBColor(0x35, 0x30, 0x20), text_color=ACCENT3, font_size=14)

add_textbox(slide, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.0),
            "현재 전략:\n"
            "• 초기: send_multimessage만 사용 (복사 가능 우선)\n"
            "• 이후: 콘텐츠 유형에 따라 두 포맷 조합\n"
            "  - 텍스트 답변 → multimessage\n"
            "  - 구조화 데이터 → richnotification\n"
            "  - 복합 → 두 포맷 함께 전송",
            font_size=14, color=WHITE)

# ══════════════════════════════════════════════════
# Slide 5: 핵심 설계 원칙 (3/3)
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "02", "핵심 설계 원칙 (3/3)")

# LLM 계층
add_accent_bar(slide, Inches(0.8), Inches(1.5), Inches(0.06), Inches(1.8), ACCENT2)
add_textbox(slide, Inches(1.1), Inches(1.45), Inches(5.3), Inches(0.5),
            "LLM 계층: LangChain 모델 레지스트리", font_size=20, color=ACCENT2, bold=True)
add_textbox(slide, Inches(1.1), Inches(2.0), Inches(5.3), Inches(1.3),
            "• 사내 LLM: OpenAI-compatible → ChatOpenAI 사용\n"
            "• api/llm/ = 모델 레지스트리 역할\n"
            "• 모델별 설정 관리 (base_url, api_key, temperature 등)\n"
            "• 기존 httpx 직접 호출 → ChatOpenAI로 대체",
            font_size=15, color=WHITE)

# 대화 히스토리
add_accent_bar(slide, Inches(6.8), Inches(1.5), Inches(0.06), Inches(1.8), ACCENT3)
add_textbox(slide, Inches(7.1), Inches(1.45), Inches(5.5), Inches(0.5),
            "대화 히스토리: Checkpointer + 아카이빙", font_size=20, color=ACCENT3, bold=True)
add_textbox(slide, Inches(7.1), Inches(2.0), Inches(5.5), Inches(1.3),
            "• MongoDBSaver가 대화 저장/복원 자동 처리\n"
            "• conversation_service.py 수동 패턴 제거\n"
            "• 완료 대화 → OpenSearch 아카이빙\n"
            "• thread_id = user_id 단독 사용",
            font_size=15, color=WHITE)

# 개발/운영 분리
add_accent_bar(slide, Inches(0.8), Inches(3.8), Inches(0.06), Inches(1.5), ORANGE)
add_textbox(slide, Inches(1.1), Inches(3.75), Inches(5.3), Inches(0.5),
            "개발/운영 환경 분리", font_size=20, color=ORANGE, bold=True)
add_textbox(slide, Inches(1.1), Inches(4.3), Inches(5.3), Inches(1.0),
            "• 개발: MemorySaver (env var로 선택)\n"
            "• 운영: MongoDBSaver\n"
            "• MongoDB 장애 시 → worker 실패 + queue 재시도 (올바른 동작)",
            font_size=15, color=WHITE)

# 운영 대시보드
add_accent_bar(slide, Inches(6.8), Inches(3.8), Inches(0.06), Inches(1.5), ORANGE)
add_textbox(slide, Inches(7.1), Inches(3.75), Inches(5.5), Inches(0.5),
            "운영 대시보드용 read model", font_size=20, color=ORANGE, bold=True)
add_textbox(slide, Inches(7.1), Inches(4.3), Inches(5.5), Inches(1.0),
            "• 단기: reply 전송 시 별도 MongoDB collection에 append\n"
            "• 중기: archive/ → OpenSearch → 대시보드 조회\n"
            "• checkpointer는 대시보드 용도에 부적합",
            font_size=15, color=WHITE)

# ══════════════════════════════════════════════════
# Slide 6: 권장 폴더 구조
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "03", "권장 api/ 폴더 구조")

# Column 1
col1 = (
    "cube/\n"
    "  client.py, router.py\n"
    "  service.py, worker.py\n"
    "  payload.py, queue.py\n"
    "  models.py\n\n"
    "chat/\n"
    "  service.py, models.py\n"
    "  history.py, commands.py"
)
add_rounded_box(slide, Inches(0.5), Inches(1.5), Inches(3.7), Inches(4.5),
                col1, fill_color=RGBColor(0x20, 0x20, 0x35), text_color=WHITE, font_size=14)

# Column 2
col2 = (
    "workflows/\n"
    "  chat/\n"
    "    graph.py, state.py\n"
    "    nodes.py, routing.py\n"
    "    prompts.py\n"
    "  rag/\n"
    "    graph.py, state.py, ...\n"
    "  agent/\n"
    "    graph.py, state.py, ..."
)
add_rounded_box(slide, Inches(4.5), Inches(1.5), Inches(3.7), Inches(4.5),
                col2, fill_color=RGBColor(0x20, 0x20, 0x35), text_color=WHITE, font_size=14)

# Column 3
col3 = (
    "llm/\n"
    "  registry.py\n"
    "  prompt/system.py\n\n"
    "mcp/\n"
    "  client.py, registry.py\n"
    "  cache.py, executor.py\n"
    "  tool_adapter.py\n"
    "  tool_selector.py"
)
add_rounded_box(slide, Inches(8.5), Inches(1.5), Inches(4.3), Inches(2.5),
                col3, fill_color=RGBColor(0x20, 0x20, 0x35), text_color=WHITE, font_size=14)

col4 = (
    "rag/\n"
    "  retriever.py, ranker.py\n"
    "  context_builder.py\n\n"
    "archive/\n"
    "  service.py, extractor.py\n"
    "  opensearch_client.py"
)
add_rounded_box(slide, Inches(8.5), Inches(4.2), Inches(4.3), Inches(2.4),
                col4, fill_color=RGBColor(0x20, 0x20, 0x35), text_color=WHITE, font_size=14)

# ══════════════════════════════════════════════════
# Slide 7: 폴더별 책임 (Cube & Chat)
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "04", "폴더별 책임 — Cube & Chat")

# Cube
add_accent_bar(slide, Inches(0.8), Inches(1.5), Inches(0.06), Inches(3.5), ACCENT)
add_textbox(slide, Inches(1.1), Inches(1.45), Inches(5.3), Inches(0.5),
            "api/cube/ — Cube 연동 경계", font_size=22, color=ACCENT, bold=True)

cube_items = [
    "HTTP endpoint 수신 & payload parsing",
    "wake-up / empty / duplicate 처리",
    "Redis queue 적재 & worker 루프 실행",
    "최종 Cube 응답 전송 및 포맷 라우팅",
    "ChatResult → multimessage / richnotification 변환",
    "LangGraph 내부 구조를 직접 알지 않음",
]
add_bullet_list(slide, Inches(1.1), Inches(2.0), Inches(5.3), Inches(3.0),
                cube_items, font_size=15, color=WHITE)

# Chat
add_accent_bar(slide, Inches(6.8), Inches(1.5), Inches(0.06), Inches(3.5), ACCENT2)
add_textbox(slide, Inches(7.1), Inches(1.45), Inches(5.5), Inches(0.5),
            "api/chat/ — 채팅 유스케이스 계층", font_size=22, color=ACCENT2, bold=True)

chat_items = [
    "run_chat_workflow() 진입점 제공",
    "요청/응답 모델 정의 (ChatResult 포함)",
    "! 접두사 command 판정 및 처리",
    "history 조회/저장 연결",
    "ChatResult 반환 (Cube 전송 안 함)",
    "LangGraph 종속 X — 유스케이스 유지",
]
add_bullet_list(slide, Inches(7.1), Inches(2.0), Inches(5.5), Inches(3.0),
                chat_items, font_size=15, color=WHITE)

# Commands box
add_rounded_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.2),
                "봇 명령어 (! 접두사):  !model [이름] — 사용할 LLM 모델 변경  |  !remove — 대화 히스토리 삭제\n"
                "command 감지 시 workflow를 타지 않고 commands.py에서 직접 처리",
                fill_color=RGBColor(0x2A, 0x2A, 0x45), text_color=ACCENT3, font_size=15)

# ══════════════════════════════════════════════════
# Slide 8: 폴더별 책임 (Workflows & LLM & MCP)
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "04", "폴더별 책임 — Workflows, LLM, MCP")

# Workflows
add_accent_bar(slide, Inches(0.5), Inches(1.5), Inches(0.06), Inches(2.5), ACCENT)
add_textbox(slide, Inches(0.8), Inches(1.45), Inches(3.7), Inches(0.5),
            "workflows/chat/", font_size=20, color=ACCENT, bold=True)
wf_items = [
    "graph 정의 & state 타입",
    "node 함수 구현",
    "conditional routing",
    "tool loop 구현",
    "fallback / 정책 분리",
]
add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(3.7), Inches(2.0),
                wf_items, font_size=14, color=WHITE)

# LLM
add_accent_bar(slide, Inches(4.6), Inches(1.5), Inches(0.06), Inches(2.5), ACCENT2)
add_textbox(slide, Inches(4.9), Inches(1.45), Inches(3.7), Inches(0.5),
            "api/llm/", font_size=20, color=ACCENT2, bold=True)
llm_items = [
    "모델 목록 관리",
    "모델별 설정 관리",
    "task 기반 모델 선택",
    "ChatOpenAI 인스턴스 생성",
    "Kimi-K2.5, Qwen3, GPT-OSS",
]
add_bullet_list(slide, Inches(4.9), Inches(2.0), Inches(3.7), Inches(2.0),
                llm_items, font_size=14, color=WHITE)

# MCP
add_accent_bar(slide, Inches(8.8), Inches(1.5), Inches(0.06), Inches(2.5), ACCENT3)
add_textbox(slide, Inches(9.1), Inches(1.45), Inches(3.7), Inches(0.5),
            "api/mcp/", font_size=20, color=ACCENT3, bold=True)
mcp_items = [
    "서버 레지스트리",
    "tool 목록 캐시",
    "MCP → OpenAI schema 변환",
    "tool name 서버 라우팅",
    "timeout / graceful degradation",
]
add_bullet_list(slide, Inches(9.1), Inches(2.0), Inches(3.7), Inches(2.0),
                mcp_items, font_size=14, color=WHITE)

# Workflow 기본 흐름
add_rounded_box(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(1.5),
                "Chat Workflow 기본 흐름:  ① context 준비 → ② 필요 tool 선택 → ③ LLM 호출 → ④ tool call 시 실행 후 재호출 → ⑤ 최종 답변 생성\n"
                "! command는 chat/service.py에서 workflow 진입 전에 처리 (workflow를 타지 않음)",
                fill_color=RGBColor(0x2A, 0x2A, 0x45), text_color=WHITE, font_size=15)

# ══════════════════════════════════════════════════
# Slide 9: 폴더별 책임 (RAG & Archive)
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "04", "폴더별 책임 — RAG & Archive")

# RAG
add_accent_bar(slide, Inches(0.8), Inches(1.5), Inches(0.06), Inches(2.0), ACCENT2)
add_textbox(slide, Inches(1.1), Inches(1.45), Inches(5.3), Inches(0.5),
            "api/rag/ — Retrieval 기능 계층", font_size=22, color=ACCENT2, bold=True)
add_textbox(slide, Inches(1.1), Inches(2.0), Inches(5.3), Inches(1.5),
            "• retriever, reranker, context builder\n"
            "• 문서 source adapter\n"
            "• 초기: chat workflow의 한 node에서 호출\n"
            "• 복잡해지면 workflows/rag/로 승격\n"
            "  (query rewrite → retrieval → rerank → compression → fallback)",
            font_size=15, color=WHITE)

# Archive
add_accent_bar(slide, Inches(6.8), Inches(1.5), Inches(0.06), Inches(2.0), ORANGE)
add_textbox(slide, Inches(7.1), Inches(1.45), Inches(5.5), Inches(0.5),
            "api/archive/ — 대화 아카이빙", font_size=22, color=ORANGE, bold=True)
add_textbox(slide, Inches(7.1), Inches(2.0), Inches(5.5), Inches(1.5),
            "• 키워드, 주제, 의도 추출 (extractor.py)\n"
            "• 사용자 만족도 신호 감지\n"
            "• OpenSearch 인덱싱 (opensearch_client.py)\n"
            "• 오래된 checkpointer 히스토리 트리밍\n"
            "• 비동기 수행 — 사용자 응답 속도 영향 없음",
            font_size=15, color=WHITE)

# 아카이빙 흐름 diagram
flow_text = (
    "아카이빙 흐름:\n"
    "cube/service.py → chat/service.py 호출 → ChatResult 수신\n"
    "                → 응답 포맷 결정 → send_multimessage / richnotification\n"
    "                → reply_sent 플래그 기록 (Redis, 중복 전송 방지)\n"
    "                → archive/service.py 호출 (비동기, 실패해도 대화 흐름 영향 없음)"
)
add_rounded_box(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(2.0),
                flow_text, fill_color=RGBColor(0x20, 0x20, 0x35), text_color=WHITE, font_size=14)

# ══════════════════════════════════════════════════
# Slide 10: 멱등성 & OpenSearch
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "05", "멱등성 전략 & OpenSearch 인덱스")

# 멱등성
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
            "멱등성(Idempotency) 전략", font_size=22, color=ACCENT, bold=True)

idem_items = [
    "Cube 응답 중복 방지: Redis에 reply_sent 플래그 (짧은 TTL)",
    "아카이빙 멱등성: OpenSearch 문서 ID = message_id (upsert)",
    "checkpointer 보호: 완료된 graph → 캐시된 결과 반환",
]
add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.0),
                idem_items, font_size=15, color=WHITE)

# 불만족 대화
add_textbox(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(0.5),
            "불만족 대화 처리", font_size=20, color=ORANGE, bold=True)
add_textbox(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(1.5),
            "• satisfaction: \"unsatisfied\" → flagged: true\n"
            "• OpenSearch에서 flagged 대화 필터링 분석\n"
            "• 프롬프트 개선, 모델 변경, RAG 소스 보강에 활용",
            font_size=15, color=WHITE)

# OpenSearch 인덱스
add_textbox(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.5),
            "OpenSearch 인덱스 구조", font_size=22, color=ACCENT2, bold=True)

os_schema = (
    '{\n'
    '  "user_id": "user-123",\n'
    '  "channel_id": "ch-456",\n'
    '  "thread_id": "user-123",\n'
    '  "timestamp": "2026-03-31T...",\n'
    '  "user_message": "...",\n'
    '  "assistant_reply": "...",\n'
    '  "keywords": ["키워드1", ...],\n'
    '  "topic": "일정 관리",\n'
    '  "satisfaction": "neutral",\n'
    '  "model_used": "kimi-k2.5",\n'
    '  "tool_calls": ["tool_name_1"],\n'
    '  "flagged": false\n'
    '}'
)
add_rounded_box(slide, Inches(6.8), Inches(2.1), Inches(5.5), Inches(4.2),
                os_schema, fill_color=RGBColor(0x20, 0x20, 0x35),
                text_color=ACCENT3, font_size=13)

# ══════════════════════════════════════════════════
# Slide 11: chat vs workflows/chat
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "06", "왜 chat과 workflows/chat을 분리하는가")

# Left: chat
add_rounded_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.2),
                "api/chat/  —  Use Case 계층\n\n"
                "• 채팅 유스케이스 계층\n"
                "• worker/상위 계층이 호출하는 진입점\n"
                "• 요청/응답 조립\n"
                "• history 저장 연결\n"
                "• ! command 판정 및 처리",
                fill_color=RGBColor(0x25, 0x35, 0x25), text_color=WHITE, font_size=16)

# Right: workflows/chat
add_rounded_box(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(3.2),
                "api/workflows/chat/  —  Engine 구현 계층\n\n"
                "• LangGraph 구현 계층\n"
                "• state, nodes, routing, tool loop\n"
                "• graph 정의 & conditional routing\n"
                "• fallback / 정책 분리\n"
                "• 재사용 가능한 LangGraph 단위",
                fill_color=RGBColor(0x25, 0x25, 0x40), text_color=WHITE, font_size=16)

# Benefits
add_textbox(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.5),
            "분리의 이점:", font_size=20, color=ACCENT, bold=True)
benefits = [
    "테스트 경계를 나누기 쉽다",
    "LangGraph 의존 범위를 좁힐 수 있다",
    "다른 orchestration 방식으로 바꿔도 외부 진입점 유지",
]
add_bullet_list(slide, Inches(0.8), Inches(5.7), Inches(11.5), Inches(1.5),
                benefits, font_size=16, color=WHITE)

# ══════════════════════════════════════════════════
# Slide 12: 마이그레이션 경로
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "07", "마이그레이션 경로")

# 현재 흐름
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.3), Inches(0.5),
            "현재 흐름", font_size=22, color=ORANGE, bold=True)
current_flow = (
    "cube/service.py:\n"
    "process_incoming_message\n\n"
    "① wake-up 판정\n"
    "② history 조회 + user message 저장\n"
    "③ generate_reply (httpx → LLM)\n"
    "④ assistant message 저장\n"
    "⑤ send_multimessage (Cube 응답)"
)
add_rounded_box(slide, Inches(0.8), Inches(2.1), Inches(5.3), Inches(3.5),
                current_flow, fill_color=RGBColor(0x35, 0x28, 0x20),
                text_color=WHITE, font_size=15)

# 변경 후 흐름
add_textbox(slide, Inches(6.6), Inches(1.5), Inches(6.2), Inches(0.5),
            "변경 후 흐름", font_size=22, color=ACCENT2, bold=True)
new_flow = (
    "cube/service.py\n"
    "  ├ wake-up / empty / duplicate 판정 (유지)\n"
    "  ├ chat/service.py 호출 → ChatResult\n"
    "  │  ├ ! command → 직접 처리\n"
    "  │  ├ workflows/chat/graph 실행\n"
    "  │  │  ├ llm/registry 모델 선택\n"
    "  │  │  ├ mcp tool calling\n"
    "  │  │  └ ChatResult 생성\n"
    "  │  └ ChatResult 반환\n"
    "  ├ 응답 포맷 결정 → 전송\n"
    "  ├ reply_sent 플래그 (Redis)\n"
    "  └ archive/service.py (비동기)"
)
add_rounded_box(slide, Inches(6.6), Inches(2.1), Inches(6.2), Inches(3.5),
                new_flow, fill_color=RGBColor(0x20, 0x30, 0x25),
                text_color=WHITE, font_size=14)

add_rounded_box(slide, Inches(0.8), Inches(6.0), Inches(12.0), Inches(0.8),
                "핵심: cube/service.py는 ChatResult를 받아 적절한 Cube 포맷으로 변환·전송하는 역할만 하고, 대화 로직은 chat/ + workflows/로 이동",
                fill_color=RGBColor(0x2A, 0x2A, 0x45), text_color=ACCENT3, font_size=15, bold=True)

# ══════════════════════════════════════════════════
# Slide 13: 초기 구현 범위
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "08", "초기 구현 범위")

# 먼저 만드는 것
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
            "먼저 만드는 것 (Phase 1)", font_size=22, color=ACCENT2, bold=True)

phase1 = [
    "api/chat/ — 채팅 유스케이스",
    "api/workflows/chat/ — LangGraph 구현",
    "api/llm/ — 레지스트리 리팩터링",
    "api/archive/ — OpenSearch 아카이빙",
    "api/mcp/ — MCP 인프라",
]
add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.5),
                phase1, font_size=17, color=WHITE)

# 나중에 만드는 것
add_textbox(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.5),
            "나중에 만드는 것 (Phase 2)", font_size=22, color=GRAY, bold=True)

phase2 = [
    "api/workflows/rag/ — RAG subgraph",
    "api/workflows/agent/ — Agent loop",
]
add_bullet_list(slide, Inches(6.8), Inches(2.1), Inches(5.5), Inches(1.5),
                phase2, font_size=17, color=GRAY)

# 전략 요약
add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.5),
            "전략:", font_size=20, color=ACCENT, bold=True)

strategy = [
    "chat workflow부터 시작",
    "MCP는 처음부터 확장형으로 분리",
    "아카이빙은 초기부터 구축하여 대화 품질 개선에 활용",
    "RAG와 Agent는 실제 복잡도가 생기면 subgraph로 승격",
]
add_bullet_list(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.0),
                strategy, font_size=17, color=WHITE)

# 필요 의존성
add_rounded_box(slide, Inches(6.8), Inches(3.8), Inches(5.5), Inches(1.5),
                "필요 의존성:\n"
                "langgraph, langchain-core, langchain-openai\n"
                "langgraph-checkpoint-mongodb, opensearch-py",
                fill_color=RGBColor(0x2A, 0x2A, 0x45), text_color=ACCENT3, font_size=14)

# ══════════════════════════════════════════════════
# Slide 14: 예상 인터페이스
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "09", "예상 인터페이스")

# chat/service.py
add_textbox(slide, Inches(0.5), Inches(1.4), Inches(3.8), Inches(0.4),
            "chat/service.py", font_size=18, color=ACCENT, bold=True)
add_rounded_box(slide, Inches(0.5), Inches(1.85), Inches(3.8), Inches(1.4),
                "def run_chat_workflow(\n"
                "    incoming: CubeIncomingMessage,\n"
                "    attempt: int = 0\n"
                ") -> ChatResult:",
                fill_color=RGBColor(0x20, 0x20, 0x35), text_color=ACCENT3, font_size=13)

# workflows/chat/state.py
add_textbox(slide, Inches(4.6), Inches(1.4), Inches(4.0), Inches(0.4),
            "workflows/chat/state.py", font_size=18, color=ACCENT, bold=True)
add_rounded_box(slide, Inches(4.6), Inches(1.85), Inches(4.0), Inches(1.4),
                "State 필드:\n"
                "messages, user_id, channel_id\n"
                "requested_model, resolved_model\n"
                "selected_tools, tool_results, error",
                fill_color=RGBColor(0x20, 0x20, 0x35), text_color=ACCENT3, font_size=13)

# llm/registry.py
add_textbox(slide, Inches(8.9), Inches(1.4), Inches(3.9), Inches(0.4),
            "llm/registry.py", font_size=18, color=ACCENT, bold=True)
add_rounded_box(slide, Inches(8.9), Inches(1.85), Inches(3.9), Inches(1.4),
                "def get_model(\n"
                "    task: str | None = None\n"
                ") -> ChatOpenAI:\n"
                "    \"\"\"task에 맞는 모델 반환\"\"\"",
                fill_color=RGBColor(0x20, 0x20, 0x35), text_color=ACCENT3, font_size=13)

# archive/service.py
add_textbox(slide, Inches(0.5), Inches(3.6), Inches(5.5), Inches(0.4),
            "archive/service.py", font_size=18, color=ACCENT, bold=True)
add_rounded_box(slide, Inches(0.5), Inches(4.05), Inches(5.5), Inches(1.8),
                "def archive_conversation(\n"
                "    thread_id: str,\n"
                "    messages: list[BaseMessage],\n"
                "    **metadata\n"
                ") -> None\n\n"
                "def trim_old_history(\n"
                "    thread_id: str, *, keep_recent: int = 20\n"
                ") -> None",
                fill_color=RGBColor(0x20, 0x20, 0x35), text_color=ACCENT3, font_size=13)

# cube/models.py
add_textbox(slide, Inches(6.3), Inches(3.6), Inches(6.2), Inches(0.4),
            "cube/models.py 확장", font_size=18, color=ACCENT, bold=True)
add_rounded_box(slide, Inches(6.3), Inches(4.05), Inches(6.2), Inches(1.0),
                "CubeIncomingMessage에 추가:\n"
                "  requested_model: str | None = None\n"
                "→ queue 통과 시에도 보존되어야 함",
                fill_color=RGBColor(0x20, 0x20, 0x35), text_color=ACCENT3, font_size=13)

# ══════════════════════════════════════════════════
# Slide 15: 테스트 계획
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "10", "테스트 계획")

tests_col1 = [
    ("test_cube_service.py", "Cube → chat 위임 검증\nempty/wake-up/duplicate 처리"),
    ("test_chat_service.py", "채팅 유스케이스 진입점\nhistory 연결 & 응답 후처리"),
    ("test_chat_graph.py", "일반 응답 / conditional routing\ntool loop 경로 검증"),
    ("test_llm_registry.py", "모델 목록 / task 기반 선택\n미등록 모델 fallback"),
]

tests_col2 = [
    ("test_mcp_registry.py", "MCP 서버 설정 로딩\nenabled/disabled 처리"),
    ("test_mcp_executor.py", "tool 서버 라우팅 / timeout\ngraceful degradation"),
    ("test_archive_service.py", "키워드·만족도 추출\nOpenSearch 인덱싱 (mock)"),
    ("test_rag_retriever.py", "retrieval & context 조합"),
]

y = Inches(1.5)
for name, desc in tests_col1:
    add_textbox(slide, Inches(0.8), y, Inches(2.5), Inches(0.4),
                name, font_size=14, color=ACCENT2, bold=True)
    add_textbox(slide, Inches(0.8), y + Inches(0.35), Inches(5.0), Inches(0.6),
                desc, font_size=12, color=GRAY)
    y += Inches(1.2)

y = Inches(1.5)
for name, desc in tests_col2:
    add_textbox(slide, Inches(6.8), y, Inches(2.5), Inches(0.4),
                name, font_size=14, color=ACCENT2, bold=True)
    add_textbox(slide, Inches(6.8), y + Inches(0.35), Inches(5.0), Inches(0.6),
                desc, font_size=12, color=GRAY)
    y += Inches(1.2)

# ══════════════════════════════════════════════════
# Slide 16: 최종 권장안
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "11", "최종 권장안")

# Summary boxes
summaries = [
    ("api/chat/", "채팅 유스케이스 (command 처리 포함)", ACCENT),
    ("api/workflows/chat/", "현재 주력 LangGraph", ACCENT),
    ("api/llm/", "LangChain 모델 레지스트리", ACCENT2),
    ("api/mcp/", "대규모 MCP tool calling 인프라", ACCENT2),
    ("api/archive/", "대화 아카이빙 (OpenSearch) + 트리밍", ORANGE),
    ("api/rag/", "retrieval 기능 계층", ORANGE),
    ("api/workflows/rag/", "나중에 RAG subgraph 시 추가", GRAY),
    ("api/workflows/agent/", "나중에 agent loop 시 추가", GRAY),
]

x_positions = [Inches(0.5), Inches(3.4), Inches(6.4), Inches(9.5)]
y_positions = [Inches(1.5), Inches(3.0)]

for idx, (name, desc, color) in enumerate(summaries):
    col = idx % 4
    row = idx // 4
    x = x_positions[col]
    y = y_positions[row]
    add_rounded_box(slide, x, y, Inches(2.7), Inches(1.2),
                    f"{name}\n\n{desc}",
                    fill_color=RGBColor(0x20, 0x20, 0x35), text_color=color, font_size=13, bold=False)

# History strategy
add_textbox(slide, Inches(0.5), Inches(4.6), Inches(12), Inches(0.5),
            "대화 히스토리 관리 전략:", font_size=20, color=ACCENT, bold=True)

hist_items = [
    "실시간: LangGraph checkpointer (MongoDBSaver) 자동 저장/복원",
    "아카이빙: 완료된 대화 → archive/ → 키워드·만족도 추출 → OpenSearch",
    "트리밍: scheduled_tasks/에서 주기적으로 오래된 히스토리 정리",
    "기존 conversation_service.py → checkpointer로 대체 후 제거",
]
add_bullet_list(slide, Inches(0.5), Inches(5.1), Inches(12), Inches(2.0),
                hist_items, font_size=16, color=WHITE)

# ══════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════
output_path = Path(__file__).parent / "LangGraph_API_구조_계획.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
